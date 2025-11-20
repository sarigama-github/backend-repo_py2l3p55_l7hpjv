import os
from io import BytesIO
import requests
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from typing import List, Optional
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class SlideItem(BaseModel):
    type: str = Field(..., description="'title', 'text', 'image', or 'equation'")
    content: Optional[str] = None
    subtitle: Optional[str] = None
    image_url: Optional[str] = None
    caption: Optional[str] = None

class SlideSpec(BaseModel):
    title: str
    subtitle: Optional[str] = None
    items: List[SlideItem] = Field(default_factory=list)

class PresentationSpec(BaseModel):
    theme_primary: str = "0b1220"  # night blue
    theme_accent: str = "d4af37"   # gold
    theme_text: str = "ffffff"
    slides: List[SlideSpec]
    filename: str = "Mesurer_la_Terre.pptx"


def hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip('#')
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


def add_title_slide(prs: Presentation, title: str, subtitle: Optional[str], colors):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    # Background color
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colors['primary']

    title_tf = slide.shapes.title.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    p.alignment = PP_ALIGN.CENTER
    run.font.name = 'Inter'
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = colors['text']

    if subtitle is not None:
        sub = slide.placeholders[1].text_frame
        sub.clear()
        p2 = sub.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        p2.alignment = PP_ALIGN.CENTER
        r2.font.name = 'Inter'
        r2.font.size = Pt(24)
        r2.font.color.rgb = colors['text']


def add_content_slide(prs: Presentation, spec: SlideSpec, colors):
    layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colors['primary']

    # Title
    title_shape = slide.shapes.title
    title_shape.text_frame.clear()
    tp = title_shape.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = spec.title
    tp.alignment = PP_ALIGN.LEFT
    tr.font.name = 'Inter'
    tr.font.size = Pt(38)
    tr.font.bold = True
    tr.font.color.rgb = colors['text']

    # Subtitle as a separate textbox if provided
    if spec.subtitle:
        box = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11), Inches(1))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = spec.subtitle
        p.alignment = PP_ALIGN.LEFT
        r.font.name = 'Inter'
        r.font.size = Pt(20)
        r.font.color.rgb = colors['text']

    y = 2.4
    for item in spec.items:
        if item.type in ['text', 'equation'] and item.content:
            box = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(7), Inches(1.2))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = item.content
            p.alignment = PP_ALIGN.LEFT
            run.font.name = 'Inter'
            run.font.size = Pt(18 if item.type == 'text' else 20)
            run.font.color.rgb = colors['text']
            if item.type == 'equation':
                run.font.bold = True
            y += 1.0
        elif item.type == 'image' and item.image_url:
            try:
                resp = requests.get(item.image_url, timeout=10)
                resp.raise_for_status()
                img = Image.open(BytesIO(resp.content))
                with BytesIO() as output:
                    img.convert('RGB').save(output, format='PNG')
                    pic_stream = output.getvalue()
                pic = slide.shapes.add_picture(BytesIO(pic_stream), Inches(8.2), Inches(y-0.2), Inches(4.5), Inches(2.8))
                if item.caption:
                    cap_box = slide.shapes.add_textbox(Inches(8.2), Inches(y+2.7), Inches(4.5), Inches(0.6))
                    tfc = cap_box.text_frame
                    tfc.clear()
                    pcap = tfc.paragraphs[0]
                    rcap = pcap.add_run()
                    rcap.text = item.caption
                    pcap.alignment = PP_ALIGN.LEFT
                    rcap.font.name = 'Inter'
                    rcap.font.size = Pt(12)
                    rcap.font.color.rgb = colors['text']
            except Exception:
                # ignore image errors, continue
                pass
            y += 3.2


def build_presentation(spec: PresentationSpec) -> bytes:
    prs = Presentation()
    # Widescreen 16:9 by default
    colors = {
        'primary': hex_to_rgb(spec.theme_primary),
        'accent': hex_to_rgb(spec.theme_accent),
        'text': hex_to_rgb(spec.theme_text)
    }

    if not spec.slides:
        raise HTTPException(status_code=400, detail="No slides provided")

    # First slide as title slide
    first = spec.slides[0]
    add_title_slide(prs, first.title, first.subtitle or "", colors)

    # Remaining slides as content
    for s in spec.slides[1:]:
        add_content_slide(prs, s, colors)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()


@app.post("/api/export/pptx")
async def export_pptx(spec: PresentationSpec):
    data = build_presentation(spec)
    from fastapi.responses import Response
    headers = {
        'Content-Disposition': f'attachment; filename="{spec.filename}"'
    }
    return Response(content=data, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers=headers)


@app.get("/")
def read_root():
    return {"message": "Backend ready for PPTX export"}


@app.get("/api/hello")
def hello():
    return {"message": "Hello from the backend API!"}


if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
